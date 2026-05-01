"""
Desktop Office Document Nodes + App Launcher + Browser
"""
import os
import subprocess
from pathlib import Path
from .base import BaseNode

import json
from typing import Any, List, Union

# Default output directory for created documents
OUTPUT_DIR = Path(os.getenv("OUTPUT_DIR", r"D:\GENAI\AI-workflow builder\outputs"))


def _ensure_output_dir():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR


def _resolve_data(params: dict, context: dict, key: str, refs: List[str]) -> Any:
    """
    Resolve data from params or context.
    1. Check the primary key (e.g. 'data', 'content').
    2. If it's a string, check if it's a reference to a context key.
    3. If still not found, check the provided ref keys (e.g. 'content_ref').
    4. Fallback to the last item in context.
    """
    val = params.get(key)

    # If key is a string reference to context
    if isinstance(val, str):
        # Strip common reference wrappers
        clean_val = val.strip("{}")
        if clean_val in context:
            val = context[clean_val]
    
    # If not found, check explicit ref keys
    if not val:
        for r in refs:
            ref_id = params.get(r)
            if isinstance(ref_id, str):
                clean_ref = ref_id.strip("{}")
                if clean_ref in context:
                    val = context[clean_ref]
                    break
                
    # Fallback to last context value if absolutely nothing specified
    if not val and context:
        val = list(context.values())[-1]
        
    return val


def _parse_to_list(data: Any) -> List[Any]:
    """Try to force data into a list format (for spreadsheets/slides)."""
    if isinstance(data, list):
        return data
    if isinstance(data, str):
        text = data.strip()
        # Try direct JSON
        try:
            parsed = json.loads(text)
            if parsed is not None:
                return parsed
        except:
            pass
        # Try Markdown JSON blocks
        import re
        match = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", text)
        if match:
            try:
                parsed = json.loads(match.group(1).strip())
                if parsed is not None:
                    return parsed
            except:
                pass
        # Fallback to CSV-like splitting
        return [line.split(",") for line in text.split("\n") if line.strip()]
    return [str(data)] if data else []


def _parse_to_str(data: Any) -> str:
    """Force data into a string format (for documents/email)."""
    if isinstance(data, str):
        return data
    if isinstance(data, (dict, list)):
        return json.dumps(data, indent=2)
    return str(data) if data is not None else ""


# ── DESKTOP_APP ─────────────────────────────────────────────────────────────

# Friendly name → executable mapping for common apps
APP_ALIASES = {
    "word":         "winword",
    "excel":        "excel",
    "powerpoint":   "powerpnt",
    "ppt":          "powerpnt",
    "notepad":      "notepad",
    "chrome":       "chrome",
    "edge":         "msedge",
    "firefox":      "firefox",
    "explorer":     "explorer",
    "calculator":   "calc",
    "paint":        "mspaint",
    "cmd":          "cmd",
    "terminal":     "wt",
    "vscode":       "code",
    "outlook":      "outlook",
    "teams":        "Teams",
}


class DesktopAppNode(BaseNode):
    node_type = "DESKTOP_APP"
    risk_level = "HIGH"

    def execute(self, params: dict, context: dict) -> str:
        app = params.get("app", "").lower().strip()
        file_path = params.get("file") or params.get("file_path")

        # Resolve alias
        exe = APP_ALIASES.get(app, app)

        from backend.config import ALLOWED_APPS
        exe_lower = exe.lower()
        if not exe_lower.endswith(".exe"):
            exe_lower += ".exe"
            
        if exe_lower not in ALLOWED_APPS:
            raise PermissionError(f"Sandboxing blocked launch of unauthorized application: '{exe_lower}'")

        # Also check context for a file path if not explicitly given
        if not file_path:
            file_ref = params.get("file_ref")
            if file_ref and file_ref in context:
                file_path = str(context[file_ref])

        try:
            if file_path:
                # Open a specific file (Windows will pick the right app)
                os.startfile(str(file_path))
                return f"Opened file '{file_path}' with the associated application."
            else:
                # Use Windows 'start' command, which searches App Paths and PATH matching the Win+R run dialog
                os.system(f'start "" "{exe}"')
                return f"Launched '{exe}' successfully."
        except Exception as e:
            return f"Failed to launch '{app}': {str(e)}"


# ── CREATE_DOCUMENT ─────────────────────────────────────────────────────────

class CreateDocumentNode(BaseNode):
    node_type = "CREATE_DOCUMENT"
    risk_level = "MEDIUM"

    def execute(self, params: dict, context: dict) -> str:
        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError:
            return "CREATE_DOCUMENT: 'python-docx' is not installed."

        _ensure_output_dir()
        filename = params.get("filename", "document.docx")
        if not filename.endswith(".docx"):
            filename += ".docx"
        filepath = OUTPUT_DIR / filename

        title = params.get("title", "AI Generated Document")
        raw_data = _resolve_data(params, context, "content", ["content_ref", "body_ref"])
        content = _parse_to_str(raw_data)

        doc = Document()

        # Title styling
        heading = doc.add_heading(title, 0)
        heading.runs[0].font.size = Pt(20)

        # Split content into paragraphs
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("# ") or line.startswith("## "):
                doc.add_heading(line.lstrip("# ").strip(), level=2)
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            else:
                doc.add_paragraph(line)

        doc.save(str(filepath))
        return str(filepath)


# ── CREATE_SPREADSHEET ──────────────────────────────────────────────────────

class CreateSpreadsheetNode(BaseNode):
    node_type = "CREATE_SPREADSHEET"
    risk_level = "MEDIUM"

    def execute(self, params: dict, context: dict) -> str:
        try:
            import openpyxl
        except ImportError:
            return "CREATE_SPREADSHEET: 'openpyxl' is not installed."

        _ensure_output_dir()
        filename = params.get("filename", "spreadsheet.xlsx")
        if not filename.endswith(".xlsx"):
            filename += ".xlsx"
        filepath = OUTPUT_DIR / filename

        sheet_name = params.get("sheet_name", "Sheet1")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name

        # Resolve and parse data
        raw_data = _resolve_data(params, context, "data", ["data_ref", "content_ref"])
        data = _parse_to_list(raw_data)

        if isinstance(data, list):
            for row in data:
                if isinstance(row, list):
                    ws.append([str(cell).strip() for cell in row])
                else:
                    ws.append([str(row)])

        wb.save(str(filepath))
        return str(filepath)


# ── CREATE_PRESENTATION ─────────────────────────────────────────────────────

class CreatePresentationNode(BaseNode):
    node_type = "CREATE_PRESENTATION"
    risk_level = "MEDIUM"

    def execute(self, params: dict, context: dict) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return "CREATE_PRESENTATION: 'python-pptx' is not installed."

        _ensure_output_dir()
        filename = params.get("filename", "presentation.pptx")
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        filepath = OUTPUT_DIR / filename

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

        # Title slide
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = params.get("title", "AI Generated Presentation")
        slide.placeholders[1].text = params.get("subtitle", "Created by AI Workflow Builder")

        # Content slides
        raw_slides = _resolve_data(params, context, "slides", ["slides_ref", "content_ref", "data_ref"])
        slides_data = _parse_to_list(raw_slides)

        # If it's just a list of strings/blobs, convert to slide dicts
        if slides_data and not isinstance(slides_data[0], dict):
            slides_data = [{"heading": f"Slide {i+1}", "body": str(item)}
                           for i, item in enumerate(slides_data)]

        for s in (slides_data or []):
            if not isinstance(s, dict):
                continue
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = str(s.get("heading", s.get("title", "")))
            slide.placeholders[1].text = str(s.get("body", s.get("content", "")))

        prs.save(str(filepath))
        return str(filepath)


# ── BROWSER_ACTION ───────────────────────────────────────────────────────────

class BrowserActionNode(BaseNode):
    node_type = "BROWSER_ACTION"
    risk_level = "HIGH"

    def execute(self, params: dict, context: dict) -> str:
        url = params.get("url") or params.get("endpoint")
        action = params.get("action", "open").lower()

        if not url:
            return "BROWSER_ACTION: 'url' is required."
        if not url.startswith("http"):
            url = "https://" + url

        if action == "open":
            import webbrowser
            webbrowser.open(url)
            return f"Opened {url} in the default browser."

        # Advanced: use playwright
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=False)
                page = browser.new_page()
                page.goto(url)

                if action == "get_text":
                    text = page.inner_text("body")
                    browser.close()
                    return text[:3000]  # cap at 3000 chars
                elif action == "screenshot":
                    shot_path = str(OUTPUT_DIR / "screenshot.png")
                    page.screenshot(path=shot_path)
                    browser.close()
                    return f"Screenshot saved: {shot_path}"
                else:
                    browser.close()
                    return f"Navigated to {url}."
        except Exception as e:
            return f"Browser automation failed: {str(e)}"
